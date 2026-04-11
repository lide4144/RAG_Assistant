from __future__ import annotations

import hashlib
import html
import json
import re
from io import StringIO
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from app.models import PageText

try:
    import fitz  # type: ignore
except ImportError:  # pragma: no cover
    fitz = None

try:
    from docx import Document as DocxDocument  # type: ignore
except ImportError:  # pragma: no cover
    DocxDocument = None

try:
    from openpyxl import load_workbook  # type: ignore
except ImportError:  # pragma: no cover
    load_workbook = None

try:
    from pptx import Presentation  # type: ignore
except ImportError:  # pragma: no cover
    Presentation = None


TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".mdx",
    ".html",
    ".htm",
    ".tex",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".csv",
    ".log",
    ".conf",
    ".ini",
    ".properties",
    ".sql",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".py",
    ".java",
    ".js",
    ".ts",
    ".swift",
    ".go",
    ".rb",
    ".php",
    ".css",
    ".scss",
    ".less",
}
OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}
LIGHT_DOCUMENT_EXTENSIONS = {".rtf", ".odt", ".epub"}
SUPPORTED_LOCAL_DOC_EXTENSIONS = {".pdf"} | TEXT_EXTENSIONS | OFFICE_EXTENSIONS | LIGHT_DOCUMENT_EXTENSIONS


@dataclass(frozen=True)
class DocumentRoute:
    source_type: str
    route_family: str
    base_parser: str
    enhanced_parser: str | None = None


def list_pdf_files(input_dir: str | Path | None) -> list[Path]:
    if input_dir is None:
        return []
    base = Path(input_dir)
    if not base.exists():
        return []
    return sorted([p for p in base.iterdir() if p.is_file() and p.suffix.lower() == ".pdf"])


def list_local_document_files(input_dir: str | Path | None) -> list[Path]:
    if input_dir is None:
        return []
    base = Path(input_dir)
    if not base.exists():
        return []
    return sorted([p for p in base.iterdir() if p.is_file() and p.suffix.lower() in SUPPORTED_LOCAL_DOC_EXTENSIONS])


def resolve_document_route(path: str | Path) -> DocumentRoute:
    suffix = Path(path).suffix.lower()
    if suffix == ".pdf":
        return DocumentRoute(source_type="pdf", route_family="pdf", base_parser="legacy", enhanced_parser="marker")
    if suffix in TEXT_EXTENSIONS:
        return DocumentRoute(source_type="text", route_family="text_like", base_parser="text")
    if suffix == ".docx":
        return DocumentRoute(source_type="docx", route_family="office", base_parser="docx")
    if suffix == ".pptx":
        return DocumentRoute(source_type="pptx", route_family="office", base_parser="pptx")
    if suffix == ".xlsx":
        return DocumentRoute(source_type="xlsx", route_family="office", base_parser="xlsx")
    if suffix in LIGHT_DOCUMENT_EXTENSIONS:
        return DocumentRoute(source_type=suffix.lstrip("."), route_family="document_like", base_parser=suffix.lstrip("."))
    return DocumentRoute(source_type="unknown", route_family="unsupported", base_parser="unsupported")


def _read_text_file(path: Path) -> str:
    payload = path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return payload.decode(encoding)
        except Exception:
            continue
    return payload.decode("utf-8", errors="ignore")


def _strip_markup(text: str) -> str:
    raw = html.unescape(str(text or ""))
    raw = re.sub(r"(?is)<script.*?>.*?</script>", " ", raw)
    raw = re.sub(r"(?is)<style.*?>.*?</style>", " ", raw)
    raw = re.sub(r"(?s)<[^>]+>", " ", raw)
    return re.sub(r"\s+", " ", raw).strip()


def _metadata_title_from_text(text: str) -> str | None:
    for line in str(text or "").splitlines():
        candidate = line.strip()
        if len(candidate) >= 4:
            return candidate[:300]
    return None


def _parse_text_like_document(path: Path) -> tuple[list[PageText], list[str], str | None]:
    suffix = path.suffix.lower()
    raw = _read_text_file(path)
    if suffix in {".html", ".htm", ".xml"}:
        text = _strip_markup(raw)
    elif suffix in {".json"}:
        try:
            text = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
        except Exception:
            text = raw
    elif suffix in {".yaml", ".yml", ".md", ".mdx", ".tex"}:
        text = raw
    elif suffix == ".csv":
        try:
            rows = list(__import__("csv").reader(StringIO(raw)))
            text = "\n".join(", ".join(cell.strip() for cell in row if cell is not None) for row in rows)
        except Exception:
            text = raw
    else:
        text = raw
    normalized = text.strip()
    if not normalized:
        return [], [f"{path.name}: empty text content"], None
    return [PageText(page_num=1, text=normalized)], [], _metadata_title_from_text(normalized)


def _parse_docx_document(path: Path) -> tuple[list[PageText], list[str], str | None]:
    if DocxDocument is None:
        return [], [f"{path.name}: python-docx is unavailable"], None
    doc = DocxDocument(str(path))
    lines: list[str] = []
    for paragraph in doc.paragraphs:
        text = str(paragraph.text or "").strip()
        if text:
            lines.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [str(cell.text or "").strip() for cell in row.cells]
            line = " | ".join(cell for cell in cells if cell)
            if line:
                lines.append(line)
    text = "\n".join(lines).strip()
    if not text:
        return [], [f"{path.name}: no readable docx text"], None
    return [PageText(page_num=1, text=text)], [], _metadata_title_from_text(text)


def _parse_pptx_document(path: Path) -> tuple[list[PageText], list[str], str | None]:
    if Presentation is None:
        return [], [f"{path.name}: python-pptx is unavailable"], None
    presentation = Presentation(str(path))
    pages: list[PageText] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                texts.append(text)
        if texts:
            pages.append(PageText(page_num=index, text="\n".join(texts)))
    if not pages:
        return [], [f"{path.name}: no readable pptx text"], None
    title = _metadata_title_from_text(pages[0].text if pages else "")
    return pages, [], title


def _parse_xlsx_document(path: Path) -> tuple[list[PageText], list[str], str | None]:
    if load_workbook is None:
        return [], [f"{path.name}: openpyxl is unavailable"], None
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    pages: list[PageText] = []
    try:
        for index, sheet in enumerate(workbook.worksheets, start=1):
            lines: list[str] = [f"# Sheet: {sheet.title}"]
            for row in sheet.iter_rows(values_only=True):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if values:
                    lines.append(", ".join(values))
            if len(lines) > 1:
                pages.append(PageText(page_num=index, text="\n".join(lines)))
    finally:
        workbook.close()
    if not pages:
        return [], [f"{path.name}: no readable xlsx text"], None
    title = _metadata_title_from_text(pages[0].text if pages else "")
    return pages, [], title


def _parse_rtf_document(path: Path) -> tuple[list[PageText], list[str], str | None]:
    raw = _read_text_file(path)
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", raw)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
    text = re.sub(r"[{}]", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return [], [f"{path.name}: no readable rtf text"], None
    return [PageText(page_num=1, text=text)], [], _metadata_title_from_text(text)


def _parse_odt_document(path: Path) -> tuple[list[PageText], list[str], str | None]:
    import zipfile

    try:
        with zipfile.ZipFile(path) as archive:
            raw = archive.read("content.xml").decode("utf-8", errors="ignore")
    except Exception as exc:
        return [], [f"{path.name}: odt extraction failed: {exc}"], None
    text = _strip_markup(raw)
    if not text:
        return [], [f"{path.name}: no readable odt text"], None
    return [PageText(page_num=1, text=text)], [], _metadata_title_from_text(text)


def _parse_epub_document(path: Path) -> tuple[list[PageText], list[str], str | None]:
    import zipfile

    try:
        with zipfile.ZipFile(path) as archive:
            html_entries = sorted(
                name
                for name in archive.namelist()
                if name.lower().endswith((".xhtml", ".html", ".htm")) and not name.endswith("/")
            )
            pages: list[PageText] = []
            for index, name in enumerate(html_entries, start=1):
                raw = archive.read(name).decode("utf-8", errors="ignore")
                text = _strip_markup(raw)
                if text:
                    pages.append(PageText(page_num=index, text=text))
    except Exception as exc:
        return [], [f"{path.name}: epub extraction failed: {exc}"], None
    if not pages:
        return [], [f"{path.name}: no readable epub text"], None
    return pages, [], _metadata_title_from_text(pages[0].text if pages else "")


def parse_local_document_pages(path: str | Path) -> tuple[list[PageText], list[str], str | None, DocumentRoute]:
    doc_path = Path(path)
    route = resolve_document_route(doc_path)
    if route.route_family == "pdf":
        pages, errors, metadata_title = parse_pdf_pages(doc_path)
        return pages, errors, metadata_title, route
    if route.route_family == "text_like":
        pages, errors, metadata_title = _parse_text_like_document(doc_path)
        return pages, errors, metadata_title, route
    if route.base_parser == "docx":
        pages, errors, metadata_title = _parse_docx_document(doc_path)
        return pages, errors, metadata_title, route
    if route.base_parser == "pptx":
        pages, errors, metadata_title = _parse_pptx_document(doc_path)
        return pages, errors, metadata_title, route
    if route.base_parser == "xlsx":
        pages, errors, metadata_title = _parse_xlsx_document(doc_path)
        return pages, errors, metadata_title, route
    if route.base_parser == "rtf":
        pages, errors, metadata_title = _parse_rtf_document(doc_path)
        return pages, errors, metadata_title, route
    if route.base_parser == "odt":
        pages, errors, metadata_title = _parse_odt_document(doc_path)
        return pages, errors, metadata_title, route
    if route.base_parser == "epub":
        pages, errors, metadata_title = _parse_epub_document(doc_path)
        return pages, errors, metadata_title, route
    return [], [f"{doc_path.name}: unsupported file type `{doc_path.suffix.lower() or 'unknown'}`"], None, route


def make_paper_id(path: Path) -> str:
    digest = hashlib.sha1(str(path.resolve()).encode("utf-8")).hexdigest()
    return digest[:12]


def pdf_content_fingerprint(path: Path) -> str:
    try:
        payload = path.read_bytes()
    except Exception:
        payload = str(path.resolve()).encode("utf-8")
    return hashlib.sha1(payload).hexdigest()


def stable_pdf_paper_id(path: Path) -> tuple[str, str]:
    fingerprint = pdf_content_fingerprint(path)
    return f"pdf_{fingerprint[:12]}", fingerprint


DEFAULT_TITLE_BLACKLIST_PATTERNS = [
    r"^\s*preprint\.?\s+under\s+review\.?\s*$",
    r"all rights reserved",
    r"copyright",
    r"arxiv preprint",
    r"provided proper attribution is provided",
    r"hereby grants permission to",
]


@dataclass
class TitleDecision:
    title: str
    source: str
    confidence: float
    adopted_layer: str = ""
    decision_trace: list[dict[str, str]] | None = None


def _normalize_structured_title_candidates(
    title_candidates: list[str] | list[dict[str, object]] | None,
) -> list[dict[str, object]]:
    normalized: list[dict[str, object]] = []
    for row in title_candidates or []:
        if isinstance(row, dict):
            text = _normalize_title_spaces(str(row.get("text", "")))
            if not text:
                continue
            normalized.append(
                {
                    "text": text[:300],
                    "source": str(row.get("source", "marker")).strip() or "marker",
                    "priority": int(row.get("priority", 99) or 99),
                    "page": int(row.get("page", 1) or 1),
                    "heading_level": row.get("heading_level"),
                    "from_markdown": bool(row.get("from_markdown")),
                }
            )
            continue
        text = _normalize_title_spaces(str(row or ""))
        if text:
            normalized.append(
                {
                    "text": text[:300],
                    "source": "marker",
                    "priority": 99,
                    "page": 1,
                    "heading_level": None,
                    "from_markdown": False,
                }
            )
    return normalized


def compile_title_blacklist_patterns(patterns: list[str] | None = None) -> list[re.Pattern[str]]:
    raw_patterns = list(patterns or DEFAULT_TITLE_BLACKLIST_PATTERNS)
    compiled: list[re.Pattern[str]] = []
    for pattern in raw_patterns:
        text = str(pattern or "").strip()
        if not text:
            continue
        try:
            compiled.append(re.compile(text, re.IGNORECASE))
        except re.error:
            continue
    if compiled:
        return compiled
    return [re.compile(p, re.IGNORECASE) for p in DEFAULT_TITLE_BLACKLIST_PATTERNS]


def _is_blacklisted_title(candidate: str, blacklist_patterns: list[re.Pattern[str]] | None = None) -> bool:
    text = str(candidate or "").strip()
    if not text:
        return True
    patterns = blacklist_patterns or compile_title_blacklist_patterns()
    for pattern in patterns:
        if pattern.search(text):
            return True
    return False


def _normalize_title_spaces(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip(" -:#\t\r\n")


def _expand_title_candidate_variants(candidate: str) -> list[str]:
    text = _normalize_title_spaces(candidate)
    if not text:
        return []

    variants: list[str] = []

    for line in re.split(r"[\r\n]+", candidate):
        normalized = _normalize_title_spaces(line)
        if len(normalized) >= 8:
            variants.append(normalized)

    without_tags = _normalize_title_spaces(re.sub(r"<[^>]+>", " ", text))
    if without_tags:
        variants.append(without_tags)

    heading_matches = re.findall(r"#\s+(.+?)(?=\s+#\s+|$)", without_tags)
    extracted_headings = [_normalize_title_spaces(match) for match in heading_matches if _normalize_title_spaces(match)]
    variants.extend(extracted_headings)
    if not extracted_headings and text:
        variants.append(text)

    cleaned_variants: list[str] = []
    seen: set[str] = set()
    for value in variants:
        normalized = _normalize_title_spaces(value)
        if not normalized:
            continue
        # Trim common author/affiliation tails after obvious separators.
        normalized = re.split(r"\babstract\b", normalized, maxsplit=1, flags=re.IGNORECASE)[0]
        normalized = re.split(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", normalized, maxsplit=1)[0]
        normalized = _normalize_title_spaces(normalized)
        if len(normalized) < 8:
            continue
        if _is_blacklisted_title(normalized):
            continue
        words = normalized.split()
        # Generate short prefixes so a good title can win over a long title+author blob.
        if len(words) >= 5:
            for size in range(5, min(len(words), 16) + 1):
                prefix = _normalize_title_spaces(" ".join(words[:size]))
                if len(prefix) >= 8 and not _is_blacklisted_title(prefix):
                    cleaned_variants.append(prefix)
        cleaned_variants.append(normalized)

    unique: list[str] = []
    for value in cleaned_variants:
        key = value.lower()
        if key in seen:
            continue
        seen.add(key)
        unique.append(value[:300])
    return unique


def score_title_candidate(candidate: str, blacklist_patterns: list[re.Pattern[str]] | None = None) -> float:
    text = str(candidate or "").strip()
    if not text:
        return 0.0
    if _is_blacklisted_title(text, blacklist_patterns=blacklist_patterns):
        return 0.0
    length = len(text)
    if length < 8:
        return 0.1
    if length > 260:
        return 0.2
    punctuation_penalty = 0.0
    if text.count(".") > 2:
        punctuation_penalty += 0.1
    if text.isupper():
        punctuation_penalty += 0.2
    words = text.split()
    if len(words) < 4:
        punctuation_penalty += 0.15
    if len(words) > 20:
        punctuation_penalty += 0.1
    if re.search(r"@[A-Za-z0-9.-]+\.", text):
        punctuation_penalty += 0.3
    if re.search(r"\b(google|university|brain|research|department)\b", text, re.IGNORECASE):
        punctuation_penalty += 0.15
    base = 0.85
    score = max(0.0, min(1.0, base - punctuation_penalty))
    return score


def choose_best_title(
    *,
    metadata_title: str | None,
    pages: list[PageText],
    title_candidates: list[str] | list[dict[str, object]] | None = None,
    confidence_threshold: float = 0.6,
    blacklist_patterns: list[str] | None = None,
) -> TitleDecision:
    compiled_blacklist = compile_title_blacklist_patterns(blacklist_patterns)
    structured_candidates = _normalize_structured_title_candidates(title_candidates)
    fallback_first_line = ""
    if pages:
        first_page = pages[0].text
        for line in first_page.splitlines():
            candidate = line.strip()
            if len(candidate) >= 8:
                fallback_first_line = candidate[:300]
                break

    layers: list[tuple[str, list[dict[str, object]]]] = [
        ("marker_h1", [row for row in structured_candidates if str(row.get("source")) == "marker_h1"]),
        ("marker_h2", [row for row in structured_candidates if str(row.get("source")) == "marker_h2"]),
        (
            "marker_markdown_first_line",
            [row for row in structured_candidates if str(row.get("source")) == "marker_markdown_first_line"],
        ),
        ("marker", [row for row in structured_candidates if str(row.get("source")) == "marker"]),
        (
            "metadata",
            [{"text": str(metadata_title).strip(), "source": "metadata", "priority": 4}] if metadata_title and str(metadata_title).strip() else [],
        ),
        (
            "fallback_first_line",
            [{"text": fallback_first_line, "source": "fallback_first_line", "priority": 5}] if fallback_first_line else [],
        ),
    ]

    trace: list[dict[str, str]] = []
    for layer_name, candidates in layers:
        if not candidates:
            trace.append({"layer": layer_name, "status": "missing", "reason": "no_candidates"})
            continue
        best_title = ""
        best_score = 0.0
        rejected_reasons: list[str] = []
        for candidate_row in candidates:
            candidate = str(candidate_row.get("text", "")).strip()
            for expanded_candidate in _expand_title_candidate_variants(candidate):
                score = score_title_candidate(expanded_candidate, blacklist_patterns=compiled_blacklist)
                if score <= 0.0:
                    rejected_reasons.append("quality_gate_rejected")
                    continue
                if score > best_score:
                    best_title = expanded_candidate[:300]
                    best_score = score
        if not best_title:
            trace.append({"layer": layer_name, "status": "rejected", "reason": rejected_reasons[-1] if rejected_reasons else "quality_gate_rejected"})
            continue
        if best_score < float(confidence_threshold):
            trace.append({"layer": layer_name, "status": "rejected", "reason": "below_confidence_threshold"})
            continue
        trace.append({"layer": layer_name, "status": "accepted", "reason": "selected"})
        return TitleDecision(
            title=best_title,
            source=layer_name,
            confidence=best_score,
            adopted_layer=layer_name,
            decision_trace=trace,
        )

        return TitleDecision(
            title="Untitled Paper",
            source="fallback_untitled",
            confidence=0.0,
            adopted_layer="fallback_untitled",
            decision_trace=trace,
        )
    return TitleDecision(
        title="Untitled Paper",
        source="fallback_untitled",
        confidence=0.0,
        adopted_layer="fallback_untitled",
        decision_trace=trace,
    )


def extract_title(metadata_title: str | None, pages: list[PageText]) -> str:
    decision = choose_best_title(
        metadata_title=metadata_title,
        pages=pages,
        title_candidates=None,
        confidence_threshold=0.0,
    )
    return decision.title


def parse_pdf_pages(pdf_path: str | Path) -> tuple[list[PageText], list[str], str | None]:
    if fitz is None:
        raise RuntimeError("PyMuPDF is required. Please install package `pymupdf`.")

    path = Path(pdf_path)
    pages: list[PageText] = []
    page_errors: list[str] = []
    metadata_title: str | None = None

    doc = fitz.open(path)
    try:
        metadata_title = (doc.metadata or {}).get("title")
        for idx in range(len(doc)):
            page_num = idx + 1
            try:
                text = doc[idx].get_text("text")
                text = text.strip()
                if text:
                    pages.append(PageText(page_num=page_num, text=text))
            except Exception as exc:  # pragma: no cover
                page_errors.append(f"{path.name}: page {page_num} parse failed: {exc}")
                continue
    finally:
        doc.close()

    return pages, page_errors, metadata_title
